
from pptx import Presentation
from pptx.util import Inches
import os

### start here

class MakeReport():    
    def __init__(self):
        # 파워포인트 객체 선언
        self.ppt = Presentation() 

    # 제목 슬라이드 추가
    def add_title(self,title_,subtitle_):
        # slide_layouts[0]: 제목 슬라이드에 해당
        title_slide = self.ppt.slide_layouts[0] 
        # 제목 슬라이드를 파워포인트 객체에 추가
        slide = self.ppt.slides.add_slide(title_slide) 

        left = Inches(0.5)
        top = Inches(0.5)
        pic = slide.shapes.add_picture("image.png",left,top)

        # 제목 - 제목에 값넣기
        title = slide.shapes.title # 제목
        title.text = title_ # 제목에 값 넣기

        # 부제목
        # 제목 상자는 placeholders[0], 부제목 상자는 [1]
        subtitle = slide.placeholders[1] 
        subtitle.text = subtitle_

        pic.zorder = 1
        title.zorder =0

        
    # 차트 슬라이드 추가  
    def add_chart1(self,text_, chart_name):
        chart_slide = self.ppt.slide_layouts[1]
        slide = self.ppt.slides.add_slide(chart_slide)

        shapes = slide.shapes
        shapes.title.text = text_
        print(shapes.title.text)

        # 차트 추가
        L = Inches(0.5)
        H = Inches(2.5)
        W = Inches(9)
        T = Inches(3)

        chart_fname = chart_name
        pic_chart = slide.shapes.add_picture(chart_fname, L, T, width=W, height=H)

    # 차트 슬라이드 추가  
    def add_chart2(self,text_, chart_name):
        chart_slide = self.ppt.slide_layouts[3]
        slide = self.ppt.slides.add_slide(chart_slide)

        shapes = slide.shapes
        shapes.title.text = text_
        print(shapes.title.text)

        # 차트 추가
        L = Inches(0.5)
        H = Inches(2.5)
        W = Inches(9)
        T = Inches(3)

        chart_fname = chart_name
        pic_chart = slide.shapes.add_picture(chart_fname, L, T, width=W, height=H)
    
    # 차트 슬라이드 추가  
    def add_chart3(self,text_, chart_name):
        chart_slide = self.ppt.slide_layouts[4]
        slide = self.ppt.slides.add_slide(chart_slide)

        shapes = slide.shapes
        shapes.title.text = text_
        print(shapes.title.text)

        # 차트 추가
        L = Inches(0.5)
        H = Inches(2.5)
        W = Inches(9)
        T = Inches(3)

        chart_fname = chart_name
        pic_chart = slide.shapes.add_picture(chart_fname, L, T, width=W, height=H)
    
    # 차트 슬라이드 추가  
    def add_chart4(self,text_, chart_name):
        chart_slide = self.ppt.slide_layouts[5]
        slide = self.ppt.slides.add_slide(chart_slide)

        shapes = slide.shapes
        shapes.title.text = text_
        print(shapes.title.text)

        # 차트 추가
        L = Inches(0.5)
        H = Inches(2.5)
        W = Inches(9)
        T = Inches(3)

        chart_fname = chart_name
        pic_chart = slide.shapes.add_picture(chart_fname, L, T, width=W, height=H)


    # 테이블 슬라이드 추가  
    def add_table_(self,text_,table_name):
        table_slide = self.ppt.slide_layouts[7]
        slide = self.ppt.slides.add_slide(table_slide)
        
        shapes = slide.shapes
        shapes.title.text = text_

        # 테이블 추가
        L = Inches(-1)
        H = Inches(3)
        W = Inches(12)
        T = Inches(2.5)

        table_fname = table_name
        pic_table = slide.shapes.add_picture(table_fname, L, T, width=W, height=H)

    #보고서 저장
    def save_report(self):
        ppt_name = os.path.join("C:/Users/user/OneDrive - pusan.ac.kr/바탕 화면/stock_report" , 'stock_report.pptx')
        try:
            # PPT 저장
            self.ppt.save(ppt_name)
            print("PPT 저장 성공")
        except Exception as e:
            # 저장 중 오류 발생 시 오류 메시지 출력
            print(f"PPT 저장 실패: {e}")
        return ppt_name
    
### end here